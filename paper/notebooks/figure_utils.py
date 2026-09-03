"""Shared figure-saving utilities for pykarambola manuscript notebooks.

Usage
-----
    from figure_utils import save_figure_pptx

    # After plt.savefig(...) or fig.savefig(...):
    save_figure_pptx('results/my_figure.pptx')          # uses current figure
    save_figure_pptx('results/my_figure.pptx', fig=fig) # explicit figure
"""

import io
from pathlib import Path


def save_figure_pptx(pptx_path, fig=None, dpi=600):
    """Save a matplotlib figure as a single-slide PPTX file.

    The slide is sized to match the figure's exact inch dimensions so the
    image fills the slide without margins, which is the standard format for
    journal figure submission.

    Parameters
    ----------
    pptx_path : str or Path
        Output path, e.g. 'results/timing_comparison_option_a.pptx'
    fig : matplotlib.figure.Figure, optional
        Figure to save.  Defaults to the current active figure (plt.gcf()).
    dpi : int
        Raster resolution for the embedded image (default 600).
    """
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Emu

    if fig is None:
        fig = plt.gcf()

    w_in, h_in = fig.get_size_inches()
    emu_per_inch = 914400  # 1 inch = 914400 EMU

    # Render figure to in-memory PNG at submission DPI
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight')
    buf.seek(0)

    prs = Presentation()
    prs.slide_width  = Emu(int(w_in * emu_per_inch))
    prs.slide_height = Emu(int(h_in * emu_per_inch))

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)

    prs.save(pptx_path)
    print(f'Saved → {pptx_path}')
